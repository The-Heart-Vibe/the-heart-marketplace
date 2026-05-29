"""
generate.py — Presentation generation engine
The Heart / VASBOX Plugin

Architecture:
- Loads a real branded template .pptx from bundled/templates/
- Clones the correct layout slide for each spec entry
- Identifies title / body / column shapes by position heuristics
- Injects text while preserving run-level font formatting
- Renders icons via cairosvg and inserts as images
- Generates charts with matplotlib and inserts as images

template_map.yaml maps each template name + layout → slide index to clone.

Usage:
    python generate.py --spec spec.json --output /path/to/output.pptx [--template blank]
"""

import argparse
import copy
import json
import os
from datetime import date
from pathlib import Path
from typing import Optional

import yaml
from pptx import Presentation
from pptx.util import Emu

# ── Paths ──────────────────────────────────────────────────────────────────────
BUNDLE_DIR    = Path(__file__).parent
ICONS_DIR     = BUNDLE_DIR / "icons"
SVG_DIR       = ICONS_DIR / "svgs"
CACHE_DIR     = ICONS_DIR / "cache"
BRAND_FILE    = BUNDLE_DIR / "brand.yaml"
TMAP_FILE     = BUNDLE_DIR / "template_map.yaml"
TEMPLATES_DIR = BUNDLE_DIR / "templates"
DEFAULT_TEMPLATE = "blank"

CACHE_DIR.mkdir(exist_ok=True)


# ── Config loaders ─────────────────────────────────────────────────────────────

def load_brand() -> dict:
    with open(BRAND_FILE) as f:
        return yaml.safe_load(f)


def load_template_map() -> dict:
    if TMAP_FILE.exists():
        with open(TMAP_FILE) as f:
            return yaml.safe_load(f) or {}
    return {}


def resolve_template(name: str) -> Path:
    """Resolve a template name to a .pptx path, with fallback to blank."""
    p = Path(name)
    if p.suffix == ".pptx" and p.is_absolute():
        if p.exists():
            return p
        raise FileNotFoundError(f"Template not found: {p}")
    stem = p.stem if p.suffix == ".pptx" else name
    candidate = TEMPLATES_DIR / f"{stem}.pptx"
    if candidate.exists():
        return candidate
    fallback = TEMPLATES_DIR / f"{DEFAULT_TEMPLATE}.pptx"
    if fallback.exists():
        print(f"  [template] '{name}' not found — falling back to {DEFAULT_TEMPLATE}.pptx")
        return fallback
    raise FileNotFoundError(f"No template found for '{name}' and no fallback at {fallback}")


def hex_to_rgb(hex_color: str) -> tuple:
    h = hex_color.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


# ── Shape detection ────────────────────────────────────────────────────────────

def is_meaningful_shape(shape, W: int, H: int) -> bool:
    """
    Filter out decorative / background shapes that shouldn't receive text.
    Removes: tiny shapes, full-slide overlays, heavily off-slide elements.
    """
    w_pct = shape.width  / W
    h_pct = shape.height / H

    # Too narrow or too short to be a content shape
    if w_pct < 0.05 or h_pct < 0.02:
        return False

    # Full-slide background overlays (tall AND wide)
    if h_pct > 0.75 and w_pct > 0.30:
        return False

    # Significantly off-slide to the left (decorative bleed elements)
    if shape.left / W < -0.10:
        return False

    return True


def get_meaningful_text_shapes(slide, W: int, H: int) -> list:
    return [
        s for s in slide.shapes
        if s.has_text_frame and is_meaningful_shape(s, W, H)
    ]


def find_shapes_for_layout(slide, layout: str, W: int, H: int) -> dict:
    """
    Return {role: shape} for the given layout using position heuristics.

    Roles by layout:
      cover     → title, subtitle
      section   → title, body (body optional / cleared)
      content   → title, body
      two_column→ title, left, right  (falls back to body if no columns found)
      chart     → title
      quote     → title, body
      closing   → cta, sub
    """
    shapes = get_meaningful_text_shapes(slide, W, H)

    def filter_content(shapes):
        """Remove obvious section-label slivers (top<5%, h<10%) and footers (top>88%)."""
        return [
            s for s in shapes
            if s.top / H < 0.88
            and not (s.top / H < 0.05 and s.height / H < 0.10)
        ]

    if layout == "cover":
        # Title = largest-area text box not in bottom 10%
        candidates = [s for s in shapes if s.top / H < 0.90]
        if not candidates:
            return {}
        title = max(candidates, key=lambda s: s.width * s.height)
        below = [s for s in candidates if s is not title and s.top > title.top]
        subtitle = min(below, key=lambda s: s.top) if below else None
        return {"title": title, "subtitle": subtitle}

    elif layout in ("section", "content", "chart", "quote"):
        candidates = filter_content(shapes)
        if not candidates:
            return {}
        title = min(candidates, key=lambda s: s.top)
        remaining = [s for s in candidates if s is not title]
        body = min(remaining, key=lambda s: s.top) if remaining else None
        return {"title": title, "body": body}

    elif layout == "two_column":
        candidates = filter_content(shapes)
        if not candidates:
            return {}
        title = min(candidates, key=lambda s: s.top)
        remaining = sorted(
            [s for s in candidates if s is not title],
            key=lambda s: s.left
        )
        result: dict = {"title": title}
        if len(remaining) >= 2:
            result["left"]  = remaining[0]
            result["right"] = remaining[-1]
        elif remaining:
            result["body"] = remaining[0]
        return result

    elif layout == "closing":
        candidates = [s for s in shapes if s.top / H < 0.85]
        if not candidates:
            return {}
        cta = max(candidates, key=lambda s: s.width * s.height)
        remaining = [s for s in candidates if s is not cta]
        sub = max(remaining, key=lambda s: s.width * s.height) if remaining else None
        return {"cta": cta, "sub": sub}

    return {}


# ── Text injection ─────────────────────────────────────────────────────────────

def _save_run_format(shape) -> dict:
    """
    Read font properties from the first text run in a shape.
    Returns a dict suitable for re-applying after tf.clear().
    """
    saved = {}
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.name:
                saved["name"] = run.font.name
            if run.font.size:
                saved["size"] = run.font.size
            if run.font.bold is not None:
                saved["bold"] = run.font.bold
            try:
                if run.font.color.type:
                    saved["color"] = run.font.color.rgb
            except Exception:
                pass
            break
        if saved:
            break
    return saved


def _apply_run_format(run, saved: dict) -> None:
    if not saved:
        return
    if "name"  in saved: run.font.name  = saved["name"]
    if "size"  in saved: run.font.size  = saved["size"]
    if "bold"  in saved: run.font.bold  = saved["bold"]
    if "color" in saved:
        try:
            run.font.color.rgb = saved["color"]
        except Exception:
            pass


def set_shape_text(shape, text: str) -> None:
    """Replace shape text with a string (supports \\n for multiple paragraphs),
    preserving font formatting."""
    if text is None:
        return
    lines = str(text).split("\n")
    fmt = _save_run_format(shape)
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if p.runs:
            _apply_run_format(p.runs[0], fmt)


def set_shape_bullets(shape, bullets: list) -> None:
    """Replace shape text with a bullet list, preserving font formatting."""
    if not bullets:
        return
    fmt = _save_run_format(shape)
    tf = shape.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        if p.runs:
            _apply_run_format(p.runs[0], fmt)


def clear_shape_text(shape) -> None:
    shape.text_frame.clear()


def clear_example_text(slide, roles: dict, W: int, H: int) -> None:
    """
    Clear example/placeholder text from all shapes on a cloned slide that are NOT
    part of the role dict (title/body/left/right/etc) and NOT page-number tokens.

    This removes leftover template example content (mission statements, stats,
    source footnotes, labels, etc.) from pitchdeck-toolkit slides.
    Branding/logo footers are image shapes, so they are unaffected.
    """
    role_shapes = [s for s in roles.values() if s is not None]
    for shape in get_meaningful_text_shapes(slide, W, H):
        if any(shape is r for r in role_shapes):
            continue
        text = shape.text_frame.text.strip()
        # Keep only page-number auto-field tokens — clear everything else
        if "‹#›" in text or "‹#›" in text:
            continue
        clear_shape_text(shape)


# ── Template map lookup ────────────────────────────────────────────────────────

def get_layout_slide_index(tname: str, layout: str, tmap: dict) -> Optional[int]:
    return (
        tmap.get("templates", {})
            .get(tname, {})
            .get("layouts", {})
            .get(layout, {})
            .get("slide")
    )


# ── Slide cloner ───────────────────────────────────────────────────────────────

def _duplicate_slide_within(prs: Presentation, slide_index: int):
    """
    Duplicate a slide WITHIN the same Presentation (same OPC package).

    Because source and destination share the same package, image/media parts
    referenced by r:embed / r:link are simply re-used — no cross-package
    blob copying is needed.  This preserves backgrounds, brand images, and
    every other visual element from the template slide.

    The entire <p:cSld> element is deep-copied so both the background (<p:bg>)
    and the shape tree (<p:spTree>) are transferred.

    Returns the new slide (appended at the end of prs.slides).
    """
    import copy as _copy
    from pptx.oxml.ns import qn

    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    src_slide    = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[6]
    new_slide    = prs.slides.add_slide(blank_layout)

    # ── Copy relationships (images, hyperlinks) ────────────────────────────
    rId_map: dict = {}
    for rId, rel in src_slide.part.rels.items():
        rtype = rel.reltype
        if "slideLayout" in rtype or "notesSlide" in rtype:
            continue
        if rel.is_external:
            new_rId = new_slide.part.relate_to(rel.target_ref, rtype, is_external=True)
        else:
            new_rId = new_slide.part.relate_to(rel.target_part, rtype)
        rId_map[rId] = new_rId

    # ── Deep-copy the entire cSld element (background + shapes) ───────────
    src_cSld = _copy.deepcopy(src_slide._element.find(qn("p:cSld")))

    # Update all r:embed / r:link references to use the new rIds
    for el in src_cSld.iter():
        for attr_key in list(el.attrib.keys()):
            if r_ns in attr_key:
                old_rId = el.get(attr_key)
                if old_rId in rId_map:
                    el.set(attr_key, rId_map[old_rId])

    # Replace new slide's cSld in-place (preserves parent element structure)
    new_cSld = new_slide._element.find(qn("p:cSld"))
    new_slide._element.replace(new_cSld, src_cSld)

    return new_slide


def _remove_slide(prs: Presentation, slide_index: int) -> None:
    """Remove the slide at slide_index from the presentation."""
    from pptx.oxml.ns import qn

    slides    = list(prs.slides)
    slide     = slides[slide_index]
    xml_slides = prs.slides._sldIdLst

    # Find the relationship id that points to this slide part
    rId = None
    for rel_id, rel in prs.part.rels.items():
        if not rel.is_external and rel.target_part is slide.part:
            rId = rel_id
            break
    if rId is None:
        return

    # Remove from the sldIdLst XML
    for sldId in list(xml_slides):
        if sldId.get(qn("r:id")) == rId:
            xml_slides.remove(sldId)
            break

    # Drop the relationship (removes the slide part from the package)
    prs.part.drop_rel(rId)


# ── Icon engine ────────────────────────────────────────────────────────────────

def _load_manifest() -> dict:
    with open(ICONS_DIR / "icon_manifest.json") as f:
        return json.load(f)


def _resolve_icon_name(keyword: str, manifest: dict) -> str:
    kw      = keyword.lower().strip()
    aliases = manifest.get("aliases", {})
    if kw in aliases:
        return aliases[kw]
    if (SVG_DIR / f"{kw}.svg").exists():
        return kw
    for alias, icon in aliases.items():
        if kw in alias or alias in kw:
            return icon
    return "star"


def get_icon_png(keyword: str, size: int = 64, color: str = "#1A1A2E") -> Optional[Path]:
    try:
        import cairosvg
    except ImportError:
        print("  [icon] cairosvg not installed — skipping icons")
        return None

    manifest  = _load_manifest()
    icon_name = _resolve_icon_name(keyword, manifest)
    svg_path  = SVG_DIR / f"{icon_name}.svg"

    if not svg_path.exists():
        print(f"  [icon] '{icon_name}' not found, skipping")
        return None

    cache_key  = f"{icon_name}_{size}_{color.lstrip('#')}"
    cache_path = CACHE_DIR / f"{cache_key}.png"

    if not cache_path.exists():
        svg_text = svg_path.read_text()
        svg_text = svg_text.replace('stroke="currentColor"', f'stroke="{color}"')
        svg_text = svg_text.replace("stroke:currentColor",  f"stroke:{color}")
        cairosvg.svg2png(
            bytestring=svg_text.encode(),
            write_to=str(cache_path),
            output_width=size,
            output_height=size,
        )
        print(f"  [icon] rendered {icon_name} @ {size}px")

    return cache_path


# ── Chart engine ───────────────────────────────────────────────────────────────

def render_chart(chart_type: str, data: dict, brand: dict, output_path: str) -> Optional[str]:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        print("  [chart] matplotlib not installed — skipping chart")
        return None

    labels    = data.get("labels", [])
    values    = data.get("values", [])
    primary   = "#" + "".join(f"{c:02X}" for c in hex_to_rgb(brand["colors"]["primary"]))
    secondary = "#" + "".join(f"{c:02X}" for c in hex_to_rgb(brand["colors"]["secondary"]))

    fig, ax = plt.subplots(figsize=(10, 5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("#F9FAFB")
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#E5E7EB")
    ax.tick_params(colors="#6B7280")

    if chart_type == "bar":
        ax.bar(labels, values, color=secondary, zorder=3)
        ax.yaxis.grid(True, color="#E5E7EB", zorder=0)
        if values:
            ax.set_ylim(0, max(values) * 1.2)
    elif chart_type == "line":
        ax.plot(labels, values, color=secondary, linewidth=2.5,
                marker="o", markerfacecolor=primary, zorder=3)
        ax.yaxis.grid(True, color="#E5E7EB", zorder=0)
    elif chart_type == "pie":
        colours = [secondary, primary, "#10B981", "#F59E0B", "#6B7280"]
        ax.pie(values, labels=labels, colors=colours[:len(values)],
               autopct="%1.0f%%", startangle=90)
        ax.axis("equal")
    else:
        ax.text(0.5, 0.5, f"Unsupported chart type: {chart_type!r}",
                ha="center", va="center", transform=ax.transAxes)

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches="tight")
    plt.close()
    return output_path


# ── Slide builder ──────────────────────────────────────────────────────────────

def build_slide(prs, slide_spec: dict, brand: dict,
                tname: str, tmap: dict, output_dir: Path,
                n_original: int) -> None:
    """
    Duplicate the appropriate template slide within prs and inject content.

    prs must be the already-opened template Presentation.
    n_original is the count of original template slides (so we clone from
    indices 0..n_original-1 only, not from slides we've already added).
    """
    layout = slide_spec.get("layout", "content")
    W, H   = prs.slide_width, prs.slide_height

    # Resolve which template slide to duplicate
    slide_index = get_layout_slide_index(tname, layout, tmap)
    if slide_index is None:
        _fallbacks = {
            "cover": 0, "section": 1, "content": 2,
            "two_column": 3, "chart": 2, "quote": 2, "closing": -1,
        }
        slide_index = _fallbacks.get(layout, 0)

    if slide_index < 0:
        slide_index = max(0, n_original + slide_index)   # e.g. -1 → last
    slide_index = min(slide_index, n_original - 1)

    new_slide = _duplicate_slide_within(prs, slide_index)
    roles     = find_shapes_for_layout(new_slide, layout, W, H)

    # Clear leftover example text from non-role shapes (e.g. pitchdeck-toolkit
    # slides carry mission statements, restaurant stats, contact details, etc.)
    clear_example_text(new_slide, roles, W, H)

    # ── Cover ──────────────────────────────────────────────────────────────────
    if layout == "cover":
        if roles.get("title"):
            set_shape_text(roles["title"], slide_spec.get("title", ""))
        if roles.get("subtitle"):
            sub = slide_spec.get("subtitle") or (
                f"{brand['identity']['company_name']}  ·  {date.today().strftime('%B %Y')}"
            )
            set_shape_text(roles["subtitle"], sub)

    # ── Section ────────────────────────────────────────────────────────────────
    elif layout == "section":
        if roles.get("title"):
            set_shape_text(roles["title"], slide_spec.get("title", ""))
        if roles.get("body"):
            if slide_spec.get("subtitle"):
                set_shape_text(roles["body"], slide_spec["subtitle"])
            else:
                clear_shape_text(roles["body"])

    # ── Content ────────────────────────────────────────────────────────────────
    elif layout == "content":
        if roles.get("title"):
            set_shape_text(roles["title"], slide_spec.get("title", ""))
        body = slide_spec.get("body", [])
        if roles.get("body"):
            if isinstance(body, list) and body:
                set_shape_bullets(roles["body"], body)
            elif isinstance(body, str) and body:
                set_shape_text(roles["body"], body)
            else:
                clear_shape_text(roles["body"])

    # ── Two-column ─────────────────────────────────────────────────────────────
    elif layout == "two_column":
        if roles.get("title"):
            set_shape_text(roles["title"], slide_spec.get("title", ""))

        left_spec  = slide_spec.get("left",  {})
        right_spec = slide_spec.get("right", {})

        def col_text(spec):
            heading = spec.get("heading", "")
            bullets = spec.get("bullets", [])
            parts = []
            if heading:
                parts.append(heading)
            parts.extend(f"• {b}" for b in bullets)
            return "\n".join(parts)

        if roles.get("left"):
            txt = col_text(left_spec)
            set_shape_text(roles["left"], txt) if txt else clear_shape_text(roles["left"])
        if roles.get("right"):
            txt = col_text(right_spec)
            set_shape_text(roles["right"], txt) if txt else clear_shape_text(roles["right"])

        # Fallback: no separate column shapes found — merge into body
        if roles.get("body") and not roles.get("left"):
            merged = col_text(left_spec) + "\n\n" + col_text(right_spec)
            set_shape_text(roles["body"], merged.strip())

    # ── Quote ──────────────────────────────────────────────────────────────────
    elif layout == "quote":
        quote_text  = slide_spec.get("quote", slide_spec.get("title", ""))
        attribution = slide_spec.get("attribution", "")
        if roles.get("title"):
            set_shape_text(roles["title"], f'"{quote_text}"')
        if roles.get("body"):
            if attribution:
                set_shape_text(roles["body"], f"— {attribution}")
            else:
                clear_shape_text(roles["body"])

    # ── Chart ──────────────────────────────────────────────────────────────────
    elif layout == "chart":
        if roles.get("title"):
            set_shape_text(roles["title"], slide_spec.get("title", ""))
        if roles.get("body"):
            clear_shape_text(roles["body"])

    # ── Closing ────────────────────────────────────────────────────────────────
    elif layout == "closing":
        cta_headline = slide_spec.get("cta_headline", "Let's talk.")
        cta_sub      = slide_spec.get("cta_sub", brand["identity"].get("website", ""))
        if roles.get("cta"):
            set_shape_text(roles["cta"], cta_headline)
        if roles.get("sub"):
            set_shape_text(roles["sub"], cta_sub)

    # ── Icon (content slides) ──────────────────────────────────────────────────
    if layout == "content" and slide_spec.get("icon"):
        icon_color = brand["colors"]["primary"]
        icon_size  = brand["icons"]["default_size"]
        icon_path  = get_icon_png(slide_spec["icon"], size=icon_size, color=icon_color)
        if icon_path and roles.get("title"):
            title_shape = roles["title"]
            icon_emu    = Emu(int(W * 0.055))
            icon_left   = title_shape.left
            icon_top    = max(0, title_shape.top - icon_emu - Emu(int(H * 0.005)))
            new_slide.shapes.add_picture(str(icon_path), icon_left, icon_top, icon_emu, icon_emu)

    # ── Chart image ────────────────────────────────────────────────────────────
    if layout == "chart" and slide_spec.get("data"):
        chart_path = str(output_dir / f"chart_{id(slide_spec)}.png")
        rendered   = render_chart(
            slide_spec.get("chart_type", "bar"),
            slide_spec["data"],
            brand,
            chart_path,
        )
        if rendered:
            if roles.get("title"):
                t = roles["title"]
                chart_top = t.top + t.height + Emu(int(H * 0.015))
            else:
                chart_top = Emu(int(H * 0.22))
            new_slide.shapes.add_picture(
                rendered,
                Emu(int(W * 0.05)),
                chart_top,
                Emu(int(W * 0.90)),
                Emu(int(H * 0.62)),
            )

    # ── Speaker notes ──────────────────────────────────────────────────────────
    if slide_spec.get("notes"):
        new_slide.notes_slide.notes_text_frame.text = slide_spec["notes"]


# ── Main generation ────────────────────────────────────────────────────────────

def generate(spec: dict, output_path: str, template_name: str = None) -> str:
    brand      = load_brand()
    tmap       = load_template_map()
    tname      = template_name or spec.get("template") or DEFAULT_TEMPLATE
    tpath      = resolve_template(tname)
    output_dir = Path(output_path).parent

    print(f"  [template] {tpath.name}")

    # Open the template as our working presentation.
    # By staying within the same OPC package, all image/media relationships
    # in the cloned slides remain valid — no cross-package copying needed.
    prs        = Presentation(str(tpath))
    n_original = len(prs.slides)

    for slide_spec in spec.get("slides", []):
        layout = slide_spec.get("layout", "content")
        label  = slide_spec.get("title") or slide_spec.get("quote") or slide_spec.get("cta_headline", "")
        print(f"  building: {layout:<12} {label[:55]}")
        build_slide(prs, slide_spec, brand, tname, tmap, output_dir, n_original)

    # Remove the original template slides (back-to-front to keep indices stable)
    for i in range(n_original - 1, -1, -1):
        _remove_slide(prs, i)

    prs.save(output_path)
    print(f"\n  done → {output_path}  ({len(prs.slides)} slides)")
    return output_path


# ── CLI ────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate a branded presentation")
    parser.add_argument("--spec",     required=True, help="Path to spec JSON")
    parser.add_argument("--output",   required=True, help="Output .pptx path")
    parser.add_argument(
        "--template", default=None,
        help=f"Template name or path (overrides spec.template). Default: {DEFAULT_TEMPLATE}"
    )
    args = parser.parse_args()

    with open(args.spec) as f:
        spec = json.load(f)

    generate(spec, args.output, template_name=args.template)
